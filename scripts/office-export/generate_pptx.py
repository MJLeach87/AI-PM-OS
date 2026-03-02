"""
PPTX generator for PM OS Office Export.
Converts a Markdown GTM/positioning artifact to a branded PowerPoint deck
using python-pptx.

Slide mapping: each H2 becomes a new slide. H1 becomes the title slide.

Branding spec: pm-os-reference/identity/BRANDING.md
Constants:      scripts/office-export/branding.py
"""

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import branding as B


# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------
SLIDE_WIDTH  = Inches(13.33)   # Widescreen 16:9
SLIDE_HEIGHT = Inches(7.5)

TITLE_BAR_HEIGHT = Inches(1.2)
CONTENT_TOP      = Inches(1.4)
CONTENT_LEFT     = Inches(0.6)
CONTENT_WIDTH    = Inches(12.0)
CONTENT_HEIGHT   = Inches(5.7)


# ---------------------------------------------------------------------------
# Color helpers
# ---------------------------------------------------------------------------

def _rgb(hex_color: str) -> RGBColor:
    r, g, b = B.hex_to_rgb(hex_color)
    return RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _set_slide_bg(slide):
    """Set slide background to SURFACE_ALT (#F9FAFB)."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(B.SURFACE_ALT)


def _add_title_bar(slide, title_text: str):
    """Add a full-width PRIMARY-coloured title bar with white text."""
    left   = Inches(0)
    top    = Inches(0)
    width  = SLIDE_WIDTH
    height = TITLE_BAR_HEIGHT

    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(B.PRIMARY)
    bar.line.fill.background()  # no border

    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.name = B.FONT_PRIMARY
    run.font.size = Pt(B.FONT_SIZE_H2)
    run.font.bold = True
    run.font.color.rgb = _rgb(B.SURFACE)


def _add_content_box(slide, content_lines: list[str]):
    """Add a text box with body content below the title bar."""
    txBox = slide.shapes.add_textbox(
        CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for line in content_lines:
        if not line.strip():
            continue

        stripped = line.strip()
        is_h3 = stripped.startswith('### ')
        is_bullet = re.match(r'^[\-\*\+]\s+', stripped) or re.match(r'^\d+\.\s+', stripped)

        # Clean text
        text = re.sub(r'^#+\s+', '', stripped)
        text = re.sub(r'^[\-\*\+]\s+', '', text)
        text = re.sub(r'^\d+\.\s+', '', text)
        text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
        text = re.sub(r'\*(.+?)\*', r'\1', text)
        text = re.sub(r'`(.+?)`', r'\1', text)
        text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT

        if is_bullet:
            p.level = 1
            text = f'• {text}'

        run = p.add_run()
        run.text = text

        if is_h3:
            run.font.name = B.FONT_PRIMARY
            run.font.size = Pt(B.FONT_SIZE_H3)
            run.font.bold = True
            run.font.color.rgb = _rgb(B.TEXT_PRIMARY)
        else:
            run.font.name = B.FONT_PRIMARY
            run.font.size = Pt(B.FONT_SIZE_BODY)
            run.font.bold = False
            run.font.color.rgb = _rgb(B.TEXT_SECONDARY)


def _add_slide_number(slide, prs: Presentation, slide_num: int):
    """Add a small slide number in the bottom-right corner."""
    left   = Inches(12.3)
    top    = Inches(7.1)
    width  = Inches(0.8)
    height = Inches(0.3)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    p = txBox.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(slide_num)
    run.font.name = B.FONT_PRIMARY
    run.font.size = Pt(B.FONT_SIZE_CAPTION)
    run.font.color.rgb = _rgb(B.TEXT_MUTED)


# ---------------------------------------------------------------------------
# Markdown section splitter
# ---------------------------------------------------------------------------

def _split_into_slides(md_text: str) -> list[dict]:
    """
    Split Markdown into slide-worthy sections.
    - H1  → title slide (index 0)
    - H2  → new slide
    - Everything else → content lines of the current slide
    Returns list of {'title': str, 'lines': [str], 'is_title_slide': bool}
    """
    slides = []
    current = None

    for line in md_text.splitlines():
        h1 = re.match(r'^#\s+(.*)', line)
        h2 = re.match(r'^##\s+(.*)', line)

        if h1:
            if current:
                slides.append(current)
            current = {'title': h1.group(1).strip(), 'lines': [], 'is_title_slide': True}
        elif h2:
            if current:
                slides.append(current)
            current = {'title': h2.group(1).strip(), 'lines': [], 'is_title_slide': False}
        else:
            if current is None:
                current = {'title': '', 'lines': [], 'is_title_slide': False}
            current['lines'].append(line)

    if current:
        slides.append(current)

    return slides


# ---------------------------------------------------------------------------
# Presentation builder
# ---------------------------------------------------------------------------

def generate(source_path: Path, out_dir: Path) -> Path:
    """Convert a Markdown file to a branded PPTX deck. Returns output path."""
    md_text = source_path.read_text(encoding='utf-8')
    slide_data = _split_into_slides(md_text)

    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # blank layout — we build everything manually

    for idx, data in enumerate(slide_data):
        if not data['title'] and not any(l.strip() for l in data['lines']):
            continue  # skip empty

        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)

        title = data['title'] or source_path.stem.replace('_', ' ')
        _add_title_bar(slide, title)

        if data['lines']:
            _add_content_box(slide, data['lines'])

        _add_slide_number(slide, prs, idx + 1)

    stem = source_path.stem
    out_path = out_dir / f'{stem}.pptx'
    prs.save(str(out_path))
    return out_path
